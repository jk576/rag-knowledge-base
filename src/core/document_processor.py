"""增强版文档处理器 - 集成 Unstructured 进行 Office 文档解析

提供高质量的文档解析，支持多种格式：
- PDF: MinerU (优先) / pypdf (备用)
- Office: Unstructured (docx, xlsx, pptx)
- 图片: pytesseract (OCR)
- 文本: 直接读取

注意：代码文件不入 RAG 索引。
代码查找应该使用 grep/ripgrep 工具，更精确、更高效。
"""

import logging
from pathlib import Path
from typing import Optional, Dict, Any, Union

from src.rag_api.config import get_settings

settings = get_settings()
logger = logging.getLogger(__name__)


class DocumentProcessor:
    """文档处理器 - 支持多种格式

    Office 文档 (docx/xlsx/pptx) 优先使用 Unstructured 解析，
    失败时回退到原生解析器。
    
    代码文件不入 RAG 索引，推荐使用 grep/ripgrep 查找代码。
    """

    def __init__(self):
        self.mineru_available = self._check_mineru()
        self.unstructured_available = self._check_unstructured()

        if self.unstructured_available:
            logger.info("✅ Unstructured 解析器已启用")
        if self.mineru_available:
            logger.info("✅ MinerU 解析器已启用")

    def _check_mineru(self) -> bool:
        """检查 MinerU 是否可用（通过 Python 3.11 子进程）"""
        try:
            import subprocess
            script_path = Path(__file__).parent.parent.parent / "scripts" / "mineru.sh"
            venv_path = Path(__file__).parent.parent.parent / ".venv-311"
            
            # 检查脚本和环境是否存在
            if not script_path.exists() or not venv_path.exists():
                return False
            
            # 尝试运行 --version 或 --help 检测可用性
            # 由于 MinerU 没有 version 参数，检查 Python 3.11 环境中的 magic_pdf
            result = subprocess.run(
                [str(venv_path / "bin" / "python"), "-c", "import magic_pdf"],
                capture_output=True,
                timeout=5
            )
            return result.returncode == 0
            
        except Exception as e:
            logger.debug(f"MinerU 检测失败: {e}")
            return False

    def _check_unstructured(self) -> bool:
        """检查 Unstructured 是否可用"""
        try:
            from unstructured.partition.docx import partition_docx
            return True
        except ImportError:
            return False

    def extract_text(self, file_path: Union[str, Path], doc_type: str) -> str:
        """提取文档文本

        Args:
            file_path: 文件路径
            doc_type: 文档类型 (pdf/docx/xlsx/pptx/image/md/txt/code)

        Returns:
            str: 提取的文本内容

        Raises:
            ValueError: 解析失败时抛出
        """
        file_path = Path(file_path)

        # Office 文档优先使用 Unstructured
        if self.unstructured_available and doc_type in ["docx", "xlsx", "pptx"]:
            try:
                return self._extract_with_unstructured(file_path, doc_type)
            except Exception as e:
                logger.warning(f"Unstructured 解析失败: {e}，回退到原生解析器")
                # 失败时回退到原生解析器
                return self._extract_with_fallback(file_path, doc_type)

        # 其他类型使用原有逻辑
        return self._extract_with_fallback(file_path, doc_type)

    def extract_structured(
        self,
        file_path: Union[str, Path],
        doc_type: str
    ) -> Dict[str, Any]:
        """提取结构化文档（返回完整结构）

        仅支持 Office 文档类型 (docx/xlsx/pptx)。

        Args:
            file_path: 文件路径
            doc_type: 文档类型

        Returns:
            Dict: 包含 text, markdown, tables, sections 等的字典

        Raises:
            ValueError: Unstructured 未安装或不支持的文档类型
        """
        file_path = Path(file_path)

        if not self.unstructured_available:
            raise ValueError("Unstructured 未安装，无法提取结构化数据")

        if doc_type not in ["docx", "xlsx", "pptx"]:
            raise ValueError(f"不支持的文档类型: {doc_type}，仅支持 docx/xlsx/pptx")

        from src.core.unstructured_parser import UnstructuredOfficeParser

        parser = UnstructuredOfficeParser()

        if doc_type == "docx":
            result = parser.parse_docx(file_path)
        elif doc_type == "xlsx":
            result = parser.parse_xlsx(file_path)
        elif doc_type == "pptx":
            result = parser.parse_pptx(file_path)
        else:
            raise ValueError(f"不支持的类型: {doc_type}")

        return {
            "text": result.text,
            "markdown": result.markdown,
            "metadata": result.metadata,
            "tables": [self._table_to_dict(t) for t in result.tables],
            "sections": [self._section_to_dict(s) for s in result.sections],
            "images": result.images,
            "page_count": result.page_count,
        }

    def _extract_with_unstructured(
        self,
        file_path: Path,
        doc_type: str
    ) -> str:
        """使用 Unstructured 提取文本

        Args:
            file_path: 文件路径
            doc_type: 文档类型

        Returns:
            str: Markdown 格式的文本
        """
        from src.core.unstructured_parser import UnstructuredOfficeParser

        parser = UnstructuredOfficeParser()

        if doc_type == "docx":
            result = parser.parse_docx(file_path)
        elif doc_type == "xlsx":
            result = parser.parse_xlsx(file_path)
        elif doc_type == "pptx":
            result = parser.parse_pptx(file_path)
        else:
            raise ValueError(f"不支持的类型: {doc_type}")

        logger.info(
            f"Unstructured 解析完成: {file_path.name} - "
            f"{len(result.tables)} 表格, {len(result.sections)} 章节"
        )

        # 返回 Markdown 格式，保留结构
        return result.markdown

    def _extract_with_fallback(self, file_path: Path, doc_type: str) -> str:
        """使用原生解析器提取（作为回退）"""
        if doc_type == "pdf":
            return self._extract_pdf(file_path)
        elif doc_type == "docx":
            return self._extract_docx(file_path)
        elif doc_type == "xlsx":
            return self._extract_xlsx(file_path)
        elif doc_type == "pptx":
            return self._extract_pptx(file_path)
        elif doc_type == "image":
            return self._extract_image(file_path)
        elif doc_type == "md":
            return self._extract_md(file_path)
        elif doc_type == "txt":
            return self._extract_txt(file_path)
        elif doc_type == "code":
            # 代码文件不入 RAG 索引
            # 推荐使用 grep/ripgrep 查找代码
            raise ValueError(f"代码文件不入 RAG 索引，请使用 grep 查找: {file_path.name}")
        else:
            return self._extract_txt(file_path)

    def _table_to_dict(self, table) -> Dict:
        """将 ParsedTable 转换为字典"""
        return {
            "caption": table.caption,
            "headers": table.headers,
            "rows": table.rows,
            "html": table.html,
        }

    def _section_to_dict(self, section) -> Dict:
        """将 ParsedSection 转换为字典"""
        return {
            "title": section.title,
            "level": section.level,
            "content": section.content,
            "start_page": section.start_page,
        }

    # ========== PDF 解析 ==========

    def _extract_pdf(self, file_path: Path) -> str:
        """提取 PDF 文本 - 优先使用 MinerU"""
        if self.mineru_available:
            return self._extract_pdf_with_mineru(file_path)
        else:
            return self._extract_pdf_with_pypdf(file_path)

    def _extract_pdf_with_mineru(self, file_path: Path) -> str:
        """使用 MinerU 提取 PDF（通过 Python 3.11 子进程）"""
        import subprocess
        import json

        try:
            script_path = Path(__file__).parent.parent.parent / "scripts" / "mineru.sh"
            result = subprocess.run(
                [str(script_path), str(file_path)],
                capture_output=True,
                text=True,
                timeout=120
            )

            if result.returncode != 0:
                raise RuntimeError(f"MinerU 执行失败: {result.stderr}")

            output = result.stdout.strip()
            lines = output.split('\n')
            for line in reversed(lines):
                line = line.strip()
                if line and line.startswith('{'):
                    try:
                        data = json.loads(line)
                        if data.get("success"):
                            return data.get("text", "")
                        else:
                            raise RuntimeError(f"MinerU 处理失败: {data.get('error')}")
                    except json.JSONDecodeError:
                        continue

            return output

        except subprocess.TimeoutExpired:
            raise ValueError("MinerU 处理超时（超过2分钟）")
        except Exception as e:
            raise ValueError(f"MinerU 调用失败: {e}")

    def _extract_pdf_with_pypdf(self, file_path: Path) -> str:
        """使用 pypdf 提取 PDF"""
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))
            text_parts = []

            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)

            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"PDF 解析失败: {e}")

    # ========== Word 解析（原生回退） ==========

    def _extract_docx(self, file_path: Path) -> str:
        """提取 Word 文档（原生解析器）"""
        try:
            from docx import Document

            doc = Document(str(file_path))
            text_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # 尝试提取表格
            for i, table in enumerate(doc.tables, 1):
                text_parts.append(f"\n[表格 {i}]")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    text_parts.append(row_text)

            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Word 文档解析失败: {e}")

    # ========== Excel 解析（原生回退） ==========

    def _extract_xlsx(self, file_path: Path) -> str:
        """提取 Excel 文本（原生解析器）"""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"# 工作表: {sheet_name}")
                text_parts.append("")

                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(
                        str(cell) for cell in row if cell is not None
                    )
                    if row_text.strip():
                        text_parts.append(row_text)

                text_parts.append("")

            return "\n".join(text_parts)
        except ImportError:
            raise ValueError("openpyxl 未安装，无法处理 Excel 文件")
        except Exception as e:
            raise ValueError(f"Excel 解析失败: {e}")

    # ========== PowerPoint 解析（原生回退） ==========

    def _extract_pptx(self, file_path: Path) -> str:
        """提取 PowerPoint 文本（原生解析器）"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"# 幻灯片 {slide_num}")
                text_parts.append("")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())

                text_parts.append("")

            return "\n".join(text_parts)
        except ImportError:
            raise ValueError("python-pptx 未安装，无法处理 PowerPoint 文件")
        except Exception as e:
            raise ValueError(f"PowerPoint 解析失败: {e}")

    # ========== 其他格式解析 ==========

    def _extract_md(self, file_path: Path) -> str:
        """提取 Markdown"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            raise ValueError(f"Markdown 读取失败: {e}")

    def _extract_txt(self, file_path: Path) -> str:
        """提取纯文本"""
        try:
            encodings = ["utf-8", "gbk", "gb2312", "latin-1"]

            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue

            raise ValueError("无法解码文件")
        except Exception as e:
            raise ValueError(f"文本读取失败: {e}")

    def _extract_image(self, file_path: Path) -> str:
        """提取图片文本 (OCR)
        
        Returns:
            提取的文本内容
            
        Raises:
            ValueError: OCR 失败或无文本内容时抛出异常
                       （不返回占位文本，避免污染搜索质量）
        """
        try:
            from PIL import Image
            import pytesseract

            # 确保 file_path 是 Path 对象
            if isinstance(file_path, str):
                file_path = Path(file_path)

            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')

            # OCR 无结果时抛出异常，不返回占位文本
            if not text.strip():
                raise ValueError(f"图片 OCR 未识别到文本内容: {file_path.name}")

            return f"[图片: {file_path.name}]\n\n{text}"

        except ImportError:
            raise ValueError("pytesseract 或 Pillow 未安装，无法处理图片")
        except Exception as e:
            # 保留原始异常信息
            if isinstance(e, ValueError):
                raise
            raise ValueError(f"图片 OCR 失败: {e}")
